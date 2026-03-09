import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PresentationGenerator:
    """
    Service responsible for compiling autonomous research insights into a polished PowerPoint presentation.
    """
    def __init__(self):
        pass

    def generate_insights_pptx(self, workspace_name: str, insights: list) -> io.BytesIO:
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Research Insights: {workspace_name}"
        subtitle.text = "Generated autonomously by ResearchHub AI"
        
        # Add a slide for each insight
        for insight in insights:
            bullet_slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            # Format title based on insight type
            insight_type_str = str(insight.get("type", "Insight")).replace("_", " ").title()
            title_shape.text = f"[{insight_type_str}] Finding"
            
            # Add insight content
            tf = body_shape.text_frame
            tf.clear() # Clear default paragraph
            
            content = insight.get("content", "")
            
            # Split by line breaks to retain formatting somewhat in bullet points
            paragraphs = content.split('\n')
            for idx, text in enumerate(paragraphs):
                if text.strip() == "":
                    continue
                p = tf.add_paragraph()
                p.text = text.strip()
                p.font.size = Pt(14)
                
                # If it looks like a bullet point from markdown (* or -), adjust level
                if text.strip().startswith("- ") or text.strip().startswith("* ") or (len(text.strip()) > 2 and text.strip()[1:3] == ". "):
                    p.level = 1
                else:
                    p.level = 0
            
        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream

presentation_generator = PresentationGenerator()
